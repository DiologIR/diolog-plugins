"""Build the Q3 engineering board pack (.pptx).

Run:  python build_deck.py
Out:  q3-engineering-board-pack.pptx  (5 slides, 16:9)
"""

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ---------------------------------------------------------------- design tokens

INK = RGBColor(0x0F, 0x17, 0x2A)
MUTED = RGBColor(0x64, 0x74, 0x8B)
FAINT = RGBColor(0x94, 0xA3, 0xB8)
ACCENT = RGBColor(0x1D, 0x4E, 0xD8)
POSITIVE = RGBColor(0x04, 0x78, 0x57)
WARN = RGBColor(0xB4, 0x53, 0x09)
RULE = RGBColor(0xE2, 0xE8, 0xF0)
PANEL = RGBColor(0xF8, 0xFA, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

FONT = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.85)
CONTENT_W = SLIDE_W - 2 * MARGIN


def text_box(slide, left, top, width, height, *, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    tf.paragraphs[0].alignment = align
    return tf


def para(tf, text, *, size, bold=False, color=INK, space_before=0,
         space_after=0, line=None, align=None, first=False, caps=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.text = text
    p.space_before = Pt(space_before)
    p.space_after = Pt(space_after)
    if line is not None:
        p.line_spacing = line
    if align is not None:
        p.alignment = align
    run = p.runs[0]
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    if caps:
        run.font._rPr.set("cap", "all")
        run.font.size = Pt(size)
    return p


def rect(slide, left, top, width, height, *, fill=None, line=None,
         line_w=Pt(1), shape=MSO_SHAPE.RECTANGLE):
    sh = slide.shapes.add_shape(shape, left, top, width, height)
    sh.shadow.inherit = False
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = line_w
    sh.text_frame.word_wrap = True
    return sh


def blank_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    return slide


def slide_header(slide, kicker, title, *, subtitle=None):
    """Standard interior-slide header: kicker, title, hairline rule."""
    tf = text_box(slide, MARGIN, Inches(0.62), CONTENT_W, Inches(0.28))
    para(tf, kicker, size=11.5, bold=True, color=ACCENT, first=True, caps=True)

    tf = text_box(slide, MARGIN, Inches(0.95), CONTENT_W, Inches(0.55))
    para(tf, title, size=30, bold=True, color=INK, first=True, line=1.0)

    top = Inches(1.62)
    if subtitle:
        tf = text_box(slide, MARGIN, Inches(1.58), CONTENT_W, Inches(0.32))
        para(tf, subtitle, size=13.5, color=MUTED, first=True, line=1.15)
        top = Inches(2.05)

    rule = rect(slide, MARGIN, top, CONTENT_W, Emu(9525), fill=RULE)
    return rule.top + rule.height


def footer(slide, page):
    tf = text_box(slide, MARGIN, Inches(6.86), Inches(6.0), Inches(0.26))
    para(tf, "Q3 FY26 Engineering Update  |  Board Pack  |  Confidential",
         size=9.5, color=FAINT, first=True)
    tf = text_box(slide, SLIDE_W - MARGIN - Inches(1.0), Inches(6.86),
                  Inches(1.0), Inches(0.26), align=PP_ALIGN.RIGHT)
    para(tf, str(page), size=9.5, color=FAINT, first=True,
         align=PP_ALIGN.RIGHT)


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def metric_tile(slide, left, top, width, height, *, label, value, unit=None,
                delta=None, delta_color=POSITIVE):
    rect(slide, left, top, width, height, fill=PANEL, line=RULE,
         line_w=Pt(0.75))
    rect(slide, left, top, Inches(0.045), height, fill=ACCENT)

    pad = Inches(0.32)
    inner_l = left + pad
    inner_w = width - pad - Inches(0.22)

    tf = text_box(slide, inner_l, top + Inches(0.30), inner_w, Inches(0.24))
    para(tf, label, size=10.5, bold=True, color=MUTED, first=True, caps=True)

    tf = text_box(slide, inner_l, top + Inches(0.62), inner_w, Inches(0.62))
    p = tf.paragraphs[0]
    p.text = value
    r = p.runs[0]
    r.font.name, r.font.size, r.font.bold = FONT, Pt(40), True
    r.font.color.rgb = INK
    if unit:
        r2 = p.add_run()
        r2.text = " " + unit
        r2.font.name, r2.font.size, r2.font.bold = FONT, Pt(16), False
        r2.font.color.rgb = MUTED

    if delta:
        tf = text_box(slide, inner_l, top + Inches(1.30), inner_w, Inches(0.26))
        para(tf, delta, size=11.5, bold=True, color=delta_color, first=True)


def bullet(tf, text, *, first=False, size=14, space_before=9, color=INK):
    p = para(tf, "—   " + text, size=size, color=color, first=first,
             space_before=space_before, line=1.28)
    return p


# ---------------------------------------------------------------- slide 1

def slide_cover(prs):
    slide = blank_slide(prs)

    rect(slide, Emu(0), Emu(0), SLIDE_W, SLIDE_H, fill=INK)
    rect(slide, Emu(0), Emu(0), Inches(0.22), SLIDE_H, fill=ACCENT)

    tf = text_box(slide, Inches(1.15), Inches(1.55), Inches(9.5), Inches(0.3))
    para(tf, "Board Pack  |  Quarter 3, FY26", size=13, bold=True,
         color=RGBColor(0x93, 0xC5, 0xFD), first=True, caps=True)

    tf = text_box(slide, Inches(1.15), Inches(2.05), Inches(10.2), Inches(1.9))
    para(tf, "Q3 Engineering Update", size=54, bold=True, color=WHITE,
         first=True, line=1.04)
    para(tf, "Auth rebuild shipped, latency down 75%, team up to 17",
         size=20, color=RGBColor(0xCB, 0xD5, 0xE1), space_before=14, line=1.2)

    rect(slide, Inches(1.15), Inches(4.42), Inches(1.4), Emu(28575),
         fill=ACCENT)

    headline = [
        ("New auth system", "Shipped in quarter"),
        ("840ms → 210ms", "p95 latency"),
        ("14 → 17", "Engineering headcount"),
        ("2 incidents", "Resolved, 4h downtime"),
    ]
    left = Inches(1.15)
    col_w = Inches(2.62)
    for i, (big, small) in enumerate(headline):
        x = left + Inches(i * 2.62)
        tf = text_box(slide, x, Inches(4.85), col_w - Inches(0.3), Inches(0.4))
        para(tf, big, size=18, bold=True, color=WHITE, first=True)
        tf = text_box(slide, x, Inches(5.28), col_w - Inches(0.3), Inches(0.3))
        para(tf, small, size=11.5, color=RGBColor(0x94, 0xA3, 0xB8),
             first=True)

    tf = text_box(slide, Inches(1.15), Inches(6.55), Inches(9.0), Inches(0.3))
    para(tf, "Prepared for the board  |  Engineering", size=11.5,
         color=RGBColor(0x64, 0x74, 0x8B), first=True)

    notes(slide, "Opening frame: a delivery quarter. One major system shipped "
                 "(auth), the platform got materially faster, the team grew, "
                 "and reliability held with two contained incidents. Q4 is a "
                 "single-focus quarter: getting off the legacy queue.")
    return slide


# ---------------------------------------------------------------- slide 2

def slide_at_a_glance(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Summary", "Q3 at a glance",
                 subtitle="Four numbers the board should take away from the "
                          "quarter.")

    top = Inches(2.42)
    h = Inches(1.78)
    gap = Inches(0.24)
    w = (CONTENT_W - 3 * gap) / 4
    tiles = [
        dict(label="Auth system", value="Shipped",
             delta="Delivered in quarter", delta_color=POSITIVE),
        dict(label="p95 latency", value="210", unit="ms",
             delta="↓ 75% from 840ms", delta_color=POSITIVE),
        dict(label="Engineering headcount", value="17",
             delta="↑ 3 from 14", delta_color=POSITIVE),
        dict(label="Incidents", value="2",
             delta="Both resolved · 4h downtime", delta_color=WARN),
    ]
    for i, t in enumerate(tiles):
        metric_tile(slide, MARGIN + i * (w + gap), top, w, h, **t)

    band_top = Inches(4.58)
    rect(slide, MARGIN, band_top, CONTENT_W, Inches(1.86), fill=PANEL,
         line=RULE, line_w=Pt(0.75))

    tf = text_box(slide, MARGIN + Inches(0.42), band_top + Inches(0.3),
                  CONTENT_W - Inches(0.84), Inches(0.26))
    para(tf, "What it means", size=10.5, bold=True, color=MUTED, first=True,
         caps=True)

    tf = text_box(slide, MARGIN + Inches(0.42), band_top + Inches(0.66),
                  CONTENT_W - Inches(0.84), Inches(1.0))
    bullet(tf, "The quarter's headline commitment — the new auth system "
               "— shipped, and the latency work landed alongside it.",
           first=True, space_before=0)
    bullet(tf, "Growth to 17 engineers is funding the Q4 legacy-queue "
               "migration without pausing feature delivery.")
    bullet(tf, "Reliability was acceptable but not free: two incidents, four "
               "hours of total downtime, both closed out.")

    footer(slide, 2)
    notes(slide, "Anchor the discussion on these four numbers. If the board "
                 "only remembers one thing, make it the latency step-change "
                 "and that auth is done. Flag the 4 hours of downtime here "
                 "rather than letting it surface later.")
    return slide


# ---------------------------------------------------------------- slide 3

def slide_delivery(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Delivery", "Shipped: the new auth system",
                 subtitle="The quarter's largest engineering commitment, "
                          "delivered and in production.")

    col_w = Inches(6.9)
    left_top = Inches(2.5)

    tf = text_box(slide, MARGIN, left_top, col_w, Inches(0.3))
    para(tf, "What landed", size=10.5, bold=True, color=MUTED, first=True,
         caps=True)

    tf = text_box(slide, MARGIN, left_top + Inches(0.42), col_w, Inches(3.1))
    bullet(tf, "The new authentication system is live and carrying production "
               "traffic.", first=True, space_before=0)
    bullet(tf, "It replaces the previous implementation end to end — no "
               "parallel legacy auth path is left running.")
    bullet(tf, "Delivered inside the quarter, alongside the latency programme "
               "rather than in place of it.")
    bullet(tf, "Gives us a single, consistent place to make future identity "
               "and access changes.")

    panel_l = MARGIN + col_w + Inches(0.45)
    panel_w = SLIDE_W - MARGIN - panel_l
    rect(slide, panel_l, left_top, panel_w, Inches(3.55), fill=PANEL,
         line=RULE, line_w=Pt(0.75))

    tf = text_box(slide, panel_l + Inches(0.4), left_top + Inches(0.34),
                  panel_w - Inches(0.8), Inches(0.26))
    para(tf, "Status", size=10.5, bold=True, color=MUTED, first=True,
         caps=True)

    tf = text_box(slide, panel_l + Inches(0.4), left_top + Inches(0.7),
                  panel_w - Inches(0.8), Inches(0.6))
    para(tf, "Shipped", size=36, bold=True, color=POSITIVE, first=True)

    rect(slide, panel_l + Inches(0.4), left_top + Inches(1.42),
         panel_w - Inches(0.8), Emu(9525), fill=RULE)

    rows = [
        ("Scope", "Full replacement"),
        ("State", "In production"),
        ("Quarter", "Q3, on plan"),
    ]
    y = left_top + Inches(1.68)
    for label, value in rows:
        tf = text_box(slide, panel_l + Inches(0.4), y,
                      Inches(1.5), Inches(0.28))
        para(tf, label, size=11.5, color=MUTED, first=True)
        tf = text_box(slide, panel_l + Inches(1.95), y,
                      panel_w - Inches(2.35), Inches(0.28),
                      align=PP_ALIGN.RIGHT)
        para(tf, value, size=11.5, bold=True, color=INK, first=True,
             align=PP_ALIGN.RIGHT)
        y += Inches(0.52)

    footer(slide, 3)
    notes(slide, "Auth was the quarter's flagship. Key point for the board: "
                 "it is a full replacement in production, not a partial "
                 "rollout, so there is no second migration hiding behind it. "
                 "Be ready for a question on what it unlocks commercially.")
    return slide


# ---------------------------------------------------------------- slide 4

def slide_performance_reliability(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Performance & reliability",
                 "Latency down 75%; two incidents contained",
                 subtitle="p95 response time fell from 840ms to 210ms. Two "
                          "incidents occurred; both are resolved.")

    top = Inches(2.5)
    chart_w = Inches(6.9)

    tf = text_box(slide, MARGIN, top, chart_w, Inches(0.3))
    para(tf, "p95 latency (ms)", size=10.5, bold=True, color=MUTED,
         first=True, caps=True)

    data = CategoryChartData()
    data.categories = ["Before (Q2 exit)", "After (Q3 exit)"]
    data.add_series("p95 latency (ms)", (840, 210))

    gf = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, MARGIN, top + Inches(0.34),
        chart_w, Inches(3.2), data)
    chart = gf.chart
    chart.has_legend = False
    chart.has_title = False

    plot = chart.plots[0]
    plot.gap_width = 120
    plot.has_data_labels = True
    labels = plot.data_labels
    labels.font.size = Pt(14)
    labels.font.bold = True
    labels.font.name = FONT
    labels.font.color.rgb = INK
    labels.position = XL_LABEL_POSITION.OUTSIDE_END

    series = plot.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = ACCENT
    pt_before = series.points[0]
    pt_before.format.fill.solid()
    pt_before.format.fill.fore_color.rgb = RGBColor(0xCB, 0xD5, 0xE1)
    pt_after = series.points[1]
    pt_after.format.fill.solid()
    pt_after.format.fill.fore_color.rgb = ACCENT

    cat_axis = chart.category_axis
    cat_axis.has_major_gridlines = False
    cat_axis.tick_labels.font.size = Pt(12)
    cat_axis.tick_labels.font.name = FONT
    cat_axis.tick_labels.font.color.rgb = MUTED
    cat_axis.format.line.color.rgb = RULE

    val_axis = chart.value_axis
    val_axis.has_major_gridlines = False
    val_axis.visible = False

    panel_l = MARGIN + chart_w + Inches(0.45)
    panel_w = SLIDE_W - MARGIN - panel_l

    metric_tile(slide, panel_l, top, panel_w, Inches(1.55),
                label="Improvement", value="75", unit="% faster",
                delta="840ms → 210ms at p95", delta_color=POSITIVE)

    inc_top = top + Inches(1.78)
    rect(slide, panel_l, inc_top, panel_w, Inches(1.76), fill=PANEL,
         line=RULE, line_w=Pt(0.75))
    rect(slide, panel_l, inc_top, Inches(0.045), Inches(1.76), fill=WARN)

    tf = text_box(slide, panel_l + Inches(0.32), inc_top + Inches(0.28),
                  panel_w - Inches(0.6), Inches(0.26))
    para(tf, "Incidents", size=10.5, bold=True, color=MUTED, first=True,
         caps=True)

    tf = text_box(slide, panel_l + Inches(0.32), inc_top + Inches(0.6),
                  panel_w - Inches(0.6), Inches(0.5))
    p = tf.paragraphs[0]
    p.text = "2"
    r = p.runs[0]
    r.font.name, r.font.size, r.font.bold = FONT, Pt(34), True
    r.font.color.rgb = INK
    r2 = p.add_run()
    r2.text = "  both resolved"
    r2.font.name, r2.font.size = FONT, Pt(13)
    r2.font.color.rgb = MUTED

    tf = text_box(slide, panel_l + Inches(0.32), inc_top + Inches(1.22),
                  panel_w - Inches(0.6), Inches(0.3))
    para(tf, "4 hours total downtime", size=12.5, bold=True, color=WARN,
         first=True)

    footer(slide, 4)
    notes(slide, "Latency is the clearest win of the quarter: p95 840ms to "
                 "210ms, a 75% reduction. Present the two incidents in the "
                 "same breath — both resolved, four hours of downtime in "
                 "aggregate across the quarter. Expect questions on root "
                 "cause and on whether the auth or latency work contributed; "
                 "have that detail ready verbally, it is deliberately not on "
                 "the slide.")
    return slide


# ---------------------------------------------------------------- slide 5

def slide_q4(prs):
    slide = blank_slide(prs)
    slide_header(slide, "Looking ahead", "Q4: migrating off the legacy queue",
                 subtitle="One priority for the quarter, resourced by the "
                          "team growth delivered in Q3.")

    top = Inches(2.5)
    left_w = Inches(7.4)

    tf = text_box(slide, MARGIN, top, left_w, Inches(0.3))
    para(tf, "The Q4 commitment", size=10.5, bold=True, color=MUTED,
         first=True, caps=True)

    tf = text_box(slide, MARGIN, top + Inches(0.42), left_w, Inches(3.1))
    bullet(tf, "Q4's engineering priority is the migration off the legacy "
               "queue.", first=True, space_before=0)
    bullet(tf, "It is the last major piece of legacy infrastructure carried "
               "into the new stack after the auth replacement.")
    bullet(tf, "The Q3 headcount increase, 14 to 17, is what makes running "
               "this alongside normal delivery viable.")
    bullet(tf, "Migration work of this shape carries downtime and cutover "
               "risk; sequencing and rollback are being planned up front.")

    panel_l = MARGIN + left_w + Inches(0.45)
    panel_w = SLIDE_W - MARGIN - panel_l

    metric_tile(slide, panel_l, top, panel_w, Inches(1.55),
                label="Team", value="17", unit="engineers",
                delta="↑ 3 in Q3 (from 14)", delta_color=POSITIVE)

    ask_top = top + Inches(1.78)
    rect(slide, panel_l, ask_top, panel_w, Inches(1.76), fill=INK)

    tf = text_box(slide, panel_l + Inches(0.36), ask_top + Inches(0.3),
                  panel_w - Inches(0.72), Inches(0.26))
    para(tf, "Focus", size=10.5, bold=True,
         color=RGBColor(0x93, 0xC5, 0xFD), first=True, caps=True)

    tf = text_box(slide, panel_l + Inches(0.36), ask_top + Inches(0.66),
                  panel_w - Inches(0.72), Inches(0.95))
    para(tf, "Legacy queue\nmigration", size=21, bold=True, color=WHITE,
         first=True, line=1.15)

    footer(slide, 5)
    notes(slide, "Close on a single, unambiguous Q4 commitment: get off the "
                 "legacy queue. Tie it back to the headcount growth so the "
                 "board sees the Q3 hiring converting into Q4 capacity. Be "
                 "explicit that migration risk is real and being planned for "
                 "— no dates or cost figures are on this slide because "
                 "none were provided.")
    return slide


# ---------------------------------------------------------------- assemble

def build(path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_cover(prs)
    slide_at_a_glance(prs)
    slide_delivery(prs)
    slide_performance_reliability(prs)
    slide_q4(prs)

    prs.save(path)
    return path


if __name__ == "__main__":
    import os
    out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                       "q3-engineering-board-pack.pptx")
    print("wrote", build(out))
